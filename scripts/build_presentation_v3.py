"""
build_presentation_v3.py
Generates Projet8_OCR_Presentation_V3.pptx — 12 slides, mentor-validated structure.
Uses V1 template's theme; builds slides from scratch with correct content.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
ROOT   = Path(r'C:\Users\Jeanne\Documents\OCR\Projet8\projet8')
CHARTS = ROOT / 'deliverables' / 'charts'
OUT    = ROOT / 'deliverables' / 'Projet8_OCR_Presentation_V3.pptx'
V1     = Path(r'C:\Users\Jeanne\Documents\OCR\Projet 08\projet8_ocr\Projet8_OCR_Presentation.pptx')

# ── Palette ────────────────────────────────────────────────────────────────────
DB  = RGBColor(0x1F, 0x4E, 0x79)   # dark blue
MB  = RGBColor(0x2E, 0x75, 0xB6)   # medium blue
LB  = RGBColor(0xBD, 0xD7, 0xEE)   # light blue
WH  = RGBColor(0xFF, 0xFF, 0xFF)   # white
BK  = RGBColor(0x26, 0x26, 0x26)   # near-black
GR  = RGBColor(0x59, 0x59, 0x59)   # gray
LGR = RGBColor(0xF2, 0xF2, 0xF2)   # light gray
BGS = RGBColor(0xF8, 0xF9, 0xFA)   # slide bg gray
LGR2= RGBColor(0xCC, 0xCC, 0xCC)   # border gray

# Slide dimensions — set from V1 in main()
W = H = None


# ── Low-level primitives ───────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])

def _r(sl, x, y, w, h, fill=None, edge=None, edge_w=Pt(0.75)):
    """Add a rectangle shape."""
    sh = sl.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if edge:
        sh.line.color.rgb = edge
        sh.line.width = edge_w
    else:
        sh.line.fill.background()
    return sh

def _t(sl, x, y, w, h, txt, sz=10, bold=False, col=BK,
       align=PP_ALIGN.LEFT, italic=False):
    """Add a single-run textbox."""
    bx = sl.shapes.add_textbox(x, y, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col
    r.font.name = 'Calibri'
    return bx

def _mt(sl, x, y, w, h, lines, sz=10, bold=False, col=BK,
        align=PP_ALIGN.LEFT, ls=1.2, sb=2):
    """Add a multi-line textbox.
    lines: list of str  OR  dict(t, sz, b, c, al, i, sb, ls)
    """
    bx = sl.shapes.add_textbox(x, y, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(ln, dict):
            p.alignment = ln.get('al', align)
            p.line_spacing = ln.get('ls', ls)
            p.space_before = Pt(ln.get('sb', sb))
            r = p.add_run(); r.text = ln.get('t', '')
            r.font.size = Pt(ln.get('sz', sz)); r.font.bold = ln.get('b', bold)
            r.font.color.rgb = ln.get('c', col); r.font.name = 'Calibri'
            r.font.italic = ln.get('i', False)
        else:
            p.alignment = align; p.line_spacing = ls; p.space_before = Pt(sb)
            r = p.add_run(); r.text = ln
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col; r.font.name = 'Calibri'
    return bx

def _img(sl, path, x, y, w, h=None):
    if h:
        sl.shapes.add_picture(str(path), x, y, w, h)
    else:
        sl.shapes.add_picture(str(path), x, y, w)

def _tbl(sl, data, x, y, w, col_ws, hfill=MB, fsz=10):
    """Add a styled table. data[0] = header row."""
    nr, nc = len(data), len(data[0])
    rh = Inches(0.32)
    tbl = sl.shapes.add_table(nr, nc, x, y, w, rh * nr).table
    for ci, cw in enumerate(col_ws):
        tbl.columns[ci].width = cw
    alts = [RGBColor(0xF0, 0xF6, 0xFF), WH]
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci); c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = hfill if ri == 0 else alts[(ri - 1) % 2]
            for p in c.text_frame.paragraphs:
                p.space_before = Pt(1); p.space_after = Pt(1)
                for r in p.runs:
                    r.font.size = Pt(fsz); r.font.name = 'Calibri'
                    r.font.bold = (ri == 0)
                    r.font.color.rgb = WH if ri == 0 else BK

def _hdr(sl, title, sub=None):
    """Standard slide header bar."""
    hh = Inches(0.78)
    _r(sl, 0, 0, W, hh, fill=DB)
    _t(sl, Inches(0.25), Inches(0.07), W - Inches(0.4), Inches(0.55),
       title, sz=18, bold=True, col=WH)
    if sub:
        _t(sl, Inches(0.25), Inches(0.56), W - Inches(0.4), Inches(0.25),
           sub, sz=9, col=LB)


# ── Slide builders ─────────────────────────────────────────────────────────────

def s01_titre(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=DB)
    _r(sl, 0, 0, W, Inches(0.1), fill=MB)
    _r(sl, 0, H - Inches(0.7), W, Inches(0.7), fill=MB)
    _mt(sl, Inches(0.6), Inches(0.9), Inches(8.8), Inches(2.6), [
        {'t': "Analyse de l'évolution du profil", 'sz': 28, 'b': True, 'c': WH,
         'al': PP_ALIGN.LEFT, 'ls': 1.08},
        {'t': 'sociodémographique des étudiants Data', 'sz': 28, 'b': True, 'c': WH,
         'al': PP_ALIGN.LEFT, 'ls': 1.08, 'sb': 0},
        {'t': 'OpenClassrooms  —  2022–2025', 'sz': 28, 'b': True, 'c': WH,
         'al': PP_ALIGN.LEFT, 'ls': 1.08, 'sb': 0},
    ])
    _t(sl, Inches(0.6), Inches(3.5), Inches(8.8), Inches(0.45),
       'Pipeline dbt · Snowflake · INSEE · COG 2024', sz=15, col=LB)
    _t(sl, Inches(0.3), H - Inches(0.58), Inches(9.4), Inches(0.45),
       'Yukel Monfray  ·  Projet 8  ·  Certification Data Analyst OCR  ·  Juin 2026',
       sz=11, col=WH)


def s02_mission(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Mission & contexte')
    Y0 = Inches(0.84)

    _t(sl, Inches(0.25), Y0, Inches(6.35), Inches(0.28),
       'LA MISSION', sz=10, bold=True, col=MB)
    _mt(sl, Inches(0.25), Y0 + Inches(0.32), Inches(6.35), Inches(2.6), [
        "● Marie-Neige (Lead Data Analyst OCR) demande une analyse rétrospective",
        "   du profil sociodémographique des étudiants parcours Data (2022–2025).",
        "",
        "● Objectifs : qualifier les données, structurer un pipeline dbt reproductible,",
        "   enrichir avec l'INSEE, livrer des tendances actionnables.",
        "",
        "● Livrables : CSV consolidé · Workflow dbt commenté · Présentation ~15 min",
    ], sz=12, col=BK, ls=1.3, sb=1)

    # RGPD box — full height of top section
    _r(sl, Inches(6.75), Y0, Inches(3.0), Inches(2.92),
       fill=LB, edge=MB)
    _t(sl, Inches(6.9), Y0 + Inches(0.1), Inches(2.75), Inches(0.32),
       'RGPD', sz=11, bold=True, col=DB)
    _mt(sl, Inches(6.9), Y0 + Inches(0.48), Inches(2.75), Inches(2.3), [
        "● USER_ID pseudonymisé",
        "   (aucune donnée nominative)",
        "",
        "● Finalité : analyse sociodémographique",
        "   agrégée uniquement",
        "",
        "● Minimisation : âge, genre, région",
        "   et année seulement",
    ], sz=11, col=DB, ls=1.25, sb=1)

    # Livrables footer
    _r(sl, 0, Inches(3.84), W, Inches(0.07), fill=MB)
    _t(sl, Inches(0.25), Inches(3.96), W - Inches(0.4), Inches(0.28),
       'LIVRABLES ATTENDUS', sz=10, bold=True, col=MB)
    _mt(sl, Inches(0.25), Inches(4.3), W - Inches(0.4), Inches(1.2), [
        {'t': 'Livrable 1  ', 'b': True, 'sz': 12, 'c': DB},
        {'t': 'CSV consolidé — mart_profil_sociodemographique.csv', 'sz': 12, 'c': BK},
        {'t': 'Livrable 2  ', 'b': True, 'sz': 12, 'c': DB, 'sb': 5},
        {'t': 'Workflow dbt commenté (GitHub + dbt Cloud job opérationnel)', 'sz': 12, 'c': BK},
        {'t': 'Livrable 3  ', 'b': True, 'sz': 12, 'c': DB, 'sb': 5},
        {'t': 'Cette présentation (12 slides · ~15 min)', 'sz': 12, 'c': BK},
    ], ls=1.2, sb=2)


def s03_sources(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Sources de données & périmètre')

    cw = Inches(3.07)
    cards = [
        ('#1  Dataset OCR',
         '4 647 lignes · 6 colonnes\nParcours Data · 2022–2025\nUSER_ID pseudonymisé'),
        ('#2  INSEE',
         'Fichier xlsx — onglets 2022–2025\nDépartement × sexe × âge quinquennal\n1 seed normalisé'),
        ('#3  COG 2024',
         '2 CSV référentiels\n101 départements · 18 régions\nMapping dept → région pour dbt'),
    ]
    for i, (title, body) in enumerate(cards):
        cx = Inches(0.22) + i * (cw + Inches(0.12))
        _r(sl, cx, Inches(0.84), cw, Inches(1.9), fill=LB, edge=MB)
        _r(sl, cx, Inches(0.84), cw, Inches(0.42), fill=MB)
        _t(sl, cx + Inches(0.1), Inches(0.88), cw - Inches(0.15), Inches(0.35),
           title, sz=12, bold=True, col=WH)
        _t(sl, cx + Inches(0.1), Inches(1.33), cw - Inches(0.15), Inches(1.3),
           body, sz=11, col=DB)

    _t(sl, Inches(0.22), Inches(2.85), W - Inches(0.35), Inches(0.28),
       'DICTIONNAIRE DES COLONNES CLÉS (dataset OCR)', sz=10, bold=True, col=MB)
    rows = [
        ['Colonne', 'Description', 'Valeurs / exemple'],
        ['USER_ID', 'Identifiant technique unique (pseudonymisé)', 'hash non nominatif'],
        ['AGE_GROUP', "Tranche d'âge au début du parcours", '25-29 ans, 30-34 ans…'],
        ['GENDER', 'Genre déclaré', "M / F / (vide → 'Non renseigné')"],
        ['REGION', 'Région de résidence', 'Île-de-France, DROM…'],
        ['YEAR_PATH_STARTED', 'Année de début du parcours', '2022 / 2023 / 2024 / 2025'],
    ]
    _tbl(sl, rows, Inches(0.22), Inches(3.18), Inches(9.56),
         [Inches(1.8), Inches(4.36), Inches(3.4)], hfill=MB, fsz=10)


def s04_methodo(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Méthodologie de collecte — du fichier brut à Snowflake')

    sw = Inches(2.3)
    steps = [
        ('01', 'Réception\ndes fichiers bruts',
         'CSV OCR · xlsx INSEE · CSV COG\ndéposés dans data/'),
        ('02', 'Conversion Python\n(dump minimal)',
         'prepare_csv_seeds.py\nprepare_insee_seed.py\nAucune logique métier'),
        ('03', 'Chargement\nSnowflake (dbt seed)',
         '4 seeds → schéma RAW\nstudents_raw · insee_population_raw\ncog_departement · cog_region'),
        ('04', 'Transformation\ndbt (ELT)',
         'dbtf build\nStaging → Intermediate → Mart\n69/69 réussis'),
    ]
    for i, (num, title, body) in enumerate(steps):
        sx = Inches(0.18) + i * (sw + Inches(0.08))
        fc = DB if i == 3 else LB
        _r(sl, sx, Inches(0.84), sw, Inches(2.75), fill=fc, edge=MB)
        _r(sl, sx + Inches(0.1), Inches(0.95), Inches(0.46), Inches(0.46),
           fill=MB if i < 3 else LB)
        _t(sl, sx + Inches(0.11), Inches(0.97), Inches(0.44), Inches(0.38),
           num, sz=13, bold=True,
           col=WH if i < 3 else DB, align=PP_ALIGN.CENTER)
        _t(sl, sx + Inches(0.1), Inches(1.55), sw - Inches(0.18), Inches(0.58),
           title, sz=11, bold=True, col=WH if i == 3 else DB)
        _t(sl, sx + Inches(0.1), Inches(2.22), sw - Inches(0.18), Inches(1.28),
           body, sz=10, col=LB if i == 3 else GR)
        if i < 3:
            _t(sl, sx + sw + Inches(0.01), Inches(2.1), Inches(0.09),
               Inches(0.4), '▶', sz=14, col=MB)

    _r(sl, 0, Inches(3.75), W, Inches(0.07), fill=MB)
    _t(sl, Inches(0.2), Inches(3.88), W - Inches(0.35), Inches(0.28),
       'POURQUOI ELT ?', sz=10, bold=True, col=MB)
    _mt(sl, Inches(0.2), Inches(4.22), W - Inches(0.35), Inches(1.25), [
        "● Python charge d'abord (Load) — puis dbt transforme dans Snowflake (Transform). Toute la logique SQL est versionnée.",
        "● Reproductible : un seul dbtf build reconstruit de bout en bout. Nouveau fichier OCR / INSEE → relancer imports + build.",
    ], sz=11, col=GR, ls=1.3, sb=4)


def s05_dag(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=BGS)
    _hdr(sl, 'Architecture du pipeline dbt')
    _img(sl, CHARTS / 'graph_dag.png', Inches(0.18), Inches(0.84), W - Inches(0.36))


def s06_enjeux(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Enjeux qualité identifiés & solutions dbt')
    rows = [
        ['Problème', 'Impact', 'Solution dans le pipeline'],
        ['Genre manquant (~27 % en 2022)',
         'Biais dans les ratios',
         "Catégorie 'Non renseigné' — pas d'imputation (RGPD)"],
        ['DROM sans équivalent INSEE\ndépartemental',
         'Pas de ratio /100k',
         "Ligne agrégat DOM → 'DROM' via CASE SQL"],
        ['Tranches âge INSEE 60-64, 65-69…\n≠ OCR (60+)',
         'Jointure impossible',
         "Rollup SQL en intermediate → '60 ans ou plus'"],
        ['Libellés âge INSEE avec underscores\net préfixe genre',
         'Matching incorrect',
         'Strip préfixe + REPLACE en staging'],
        ['Double comptage DOM\n(depts 971-976 + agrégat)',
         'Surreprésentation DROM',
         "Depts 971-976 exclus — seul agrégat 'DOM' conservé"],
        ['USER_ID longitudinal\n(même étudiant, plusieurs années)',
         'Faux doublons',
         'Données longitudinales documentées + test singulier'],
    ]
    _tbl(sl, rows, Inches(0.22), Inches(0.85), Inches(9.56),
         [Inches(2.3), Inches(1.86), Inches(5.4)], hfill=DB, fsz=9)


def s07_tests(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Tests dbt & garantie de reproductibilité')

    _t(sl, Inches(0.25), Inches(0.84), Inches(5.0), Inches(0.3),
       'STRATÉGIE PAR COUCHE', sz=10, bold=True, col=MB)
    layers = [
        ('Sources (RAW)', 'not_null sur clé primaire — sanity check'),
        ('Staging', 'not_null · unique · accepted_values · relationships'),
        ('Intermediate', 'not_null · unique · accepted_values · relationships'),
        ('Mart', 'not_null · accepted_values'),
        ('Test singulier', 'assert_unique_student_year — 1 étudiant max / année'),
    ]
    for i, (lay, desc) in enumerate(layers):
        ly = Inches(1.2) + i * Inches(0.74)
        _r(sl, Inches(0.25), ly, Inches(1.7), Inches(0.6), fill=LB, edge=MB)
        _t(sl, Inches(0.3), ly + Inches(0.1), Inches(1.62), Inches(0.44),
           lay, sz=9.5, bold=True, col=DB)
        _t(sl, Inches(2.08), ly + Inches(0.12), Inches(3.1), Inches(0.44),
           desc, sz=10, col=BK)

    _r(sl, Inches(5.35), Inches(0.8), Inches(0.04), Inches(4.7), fill=LB)

    _t(sl, Inches(5.55), Inches(0.84), Inches(4.2), Inches(0.3),
       'REPRODUCTIBILITÉ', sz=10, bold=True, col=MB)
    guarantees = [
        ('Si : Nouvelle région apparaît', '→ accepted_values échoue · pipeline bloqué'),
        ('Si : USER_ID dupliqué dans une année', '→ test singulier échoue · pipeline bloqué'),
        ('Si : Jointure COG rompue', '→ relationships échoue · pipeline bloqué'),
        ('Si : Nouveau fichier OCR / INSEE', '→ run_all_imports + dbtf build · rebuild complet'),
    ]
    for i, (trigger, effect) in enumerate(guarantees):
        gy = Inches(1.2) + i * Inches(0.9)
        _r(sl, Inches(5.55), gy, Inches(4.2), Inches(0.76), fill=LGR, edge=LGR2)
        _t(sl, Inches(5.65), gy + Inches(0.06), Inches(4.0), Inches(0.3),
           trigger, sz=10, bold=True, col=DB)
        _t(sl, Inches(5.65), gy + Inches(0.4), Inches(4.0), Inches(0.3),
           effect, sz=10, col=BK)

    _r(sl, Inches(0.25), Inches(4.9), Inches(9.5), Inches(0.58), fill=DB)
    _t(sl, Inches(0.35), Inches(5.03), Inches(9.3), Inches(0.42),
       'dbtf build  ·  seeds + modèles + tests en une commande  ·  69 / 69 réussis',
       sz=13, bold=True, col=WH, align=PP_ALIGN.CENTER)


def s08_regions(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Résultats — Concentration régionale & évolution des inscriptions')
    _img(sl, CHARTS / 'graph_regions.png', Inches(0.18), Inches(0.84), Inches(5.9))
    _img(sl, CHARTS / 'graph_yearly.png',  Inches(6.2),  Inches(0.84), Inches(3.7))
    _r(sl, Inches(0.18), Inches(4.42), Inches(9.6), Inches(0.06), fill=MB)
    _mt(sl, Inches(0.18), Inches(4.54), Inches(9.6), Inches(1.0), [
        {'t': '▶ IDF : 45,6 % des inscriptions cumulées 2022–2025.',
         'b': True, 'sz': 12, 'c': DB},
        {'t': '▶ Baisse 2022→2024 (−50 %), rebond 2025 (+12 %).   '
              '▶ DROM : 46 inscrits (1 %) — pas de benchmark INSEE.',
         'sz': 12, 'c': BK, 'sb': 5},
    ], ls=1.2)


def s09_genre(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Résultats — Genre & qualité de la donnée')
    _img(sl, CHARTS / 'graph_nr_rate.png', Inches(0.18), Inches(0.84), Inches(5.7))

    _t(sl, Inches(6.15), Inches(0.9), Inches(3.6), Inches(0.32),
       'LECTURE', sz=10, bold=True, col=MB)
    _mt(sl, Inches(6.15), Inches(1.27), Inches(3.6), Inches(2.5), [
        "● Amélioration spectaculaire de la qualité de collecte du genre sur 4 ans.",
        "",
        "● 41,6 % → 6,7 % : cette tendance est un indicateur de pilotage qualité en soi.",
        "",
        "● À maintenir : incitations à la saisie lors de l'inscription.",
    ], sz=12, col=BK, ls=1.3, sb=2)

    _r(sl, Inches(6.15), Inches(3.9), Inches(3.6), Inches(1.55), fill=LB, edge=MB)
    _t(sl, Inches(6.28), Inches(4.0), Inches(3.35), Inches(0.32),
       'Répartition avec genre renseigné', sz=10, bold=True, col=DB)
    _mt(sl, Inches(6.28), Inches(4.38), Inches(3.35), Inches(1.0), [
        '● Hommes : ~69 % des inscrits renseignés',
        '● Femmes : ~31 % — sous-représentation persistante',
    ], sz=11, col=DB, ls=1.35, sb=2)


def s10_insee(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, "Résultats — Taux d'inscription pour 100 000 habitants")
    _img(sl, CHARTS / 'graph_100k.png', Inches(0.18), Inches(0.84), Inches(6.9))

    _t(sl, Inches(7.28), Inches(0.9), Inches(2.5), Inches(0.32),
       'LECTURE', sz=10, bold=True, col=MB)
    _mt(sl, Inches(7.28), Inches(1.27), Inches(2.5), Inches(4.0), [
        "● IDF domine encore plus après normalisation (~670 / 100k).",
        "",
        "● ARA, 2ème en effectifs bruts, descend à la 7ème place.",
        "",
        "● PACA et Centre-Val de Loire dépassent leur part d'effectifs.",
        "",
        "● Les effectifs bruts seuls sont trompeurs.",
        "",
        "● DROM exclu : pas de données INSEE métropole.",
    ], sz=11, col=BK, ls=1.3, sb=3)


def s11_reco(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Recommandations')
    recs = [
        ('01', 'Creuser la concentration IDF',
         "Analyser via le taux /100k, pas seulement les effectifs bruts. "
         "Identifier si c'est une sur-représentation réelle ou un effet de taille."),
        ('02', 'Maintenir les incitations à la saisie du genre',
         "La tendance 41 % → 7 % est très positive. "
         "Maintenir les dispositifs qui ont produit cette amélioration."),
        ('03', 'Traiter DROM comme segment à part',
         "46 inscrits (1 %) — pas de benchmark INSEE disponible. "
         "Nécessite une collecte spécifique pour un ratio fiable."),
        ('04', 'Industrialiser le rebuild annuel',
         "À chaque nouveau fichier OCR ou mise à jour INSEE : "
         "run_all_imports.py + dbtf build. Le pipeline est prêt."),
    ]
    rw = Inches(4.42); rh = Inches(1.92)
    for i, (num, title, body) in enumerate(recs):
        rx = Inches(0.22) + (i % 2) * (rw + Inches(0.2))
        ry = Inches(0.9) + (i // 2) * (rh + Inches(0.18))
        _r(sl, rx, ry, rw, rh, fill=LGR, edge=MB)
        _r(sl, rx, ry, Inches(0.48), rh, fill=MB)
        _t(sl, rx + Inches(0.06), ry + rh / 2 - Inches(0.2), Inches(0.38),
           Inches(0.38), num, sz=13, bold=True, col=WH, align=PP_ALIGN.CENTER)
        _t(sl, rx + Inches(0.56), ry + Inches(0.18), rw - Inches(0.65),
           Inches(0.36), title, sz=10, bold=True, col=DB)
        _t(sl, rx + Inches(0.56), ry + Inches(0.6), rw - Inches(0.65),
           Inches(1.22), body, sz=9, col=BK)


def s12_livrables(prs):
    sl = _blank(prs)
    _r(sl, 0, 0, W, H, fill=WH)
    _hdr(sl, 'Livrables produits & questions')

    delivs = [
        ('1', 'CSV consolidé',
         'mart_profil_sociodemographique.csv\nGrain : année × région × genre × âge\n1 741 lignes · 7 colonnes'),
        ('2', 'Workflow dbt commenté',
         'GitHub + dbt Cloud opérationnel\nSQL commenté · tests justifiés\nDoc blocks dbt'),
        ('3', 'Cette présentation',
         '12 slides · plan mentor-validé\nGraphiques depuis le mart\nReproductible'),
    ]
    dw = Inches(3.0)
    for i, (num, title, body) in enumerate(delivs):
        dx = Inches(0.22) + i * (dw + Inches(0.15))
        _r(sl, dx, Inches(0.9), dw, Inches(1.72), fill=DB)
        _r(sl, dx, Inches(0.9), dw, Inches(0.44), fill=MB)
        _t(sl, dx + Inches(0.1), Inches(0.94), Inches(0.44), Inches(0.36),
           num, sz=14, bold=True, col=WH, align=PP_ALIGN.CENTER)
        _t(sl, dx + Inches(0.58), Inches(0.96), dw - Inches(0.7),
           Inches(0.36), title, sz=10, bold=True, col=WH)
        _t(sl, dx + Inches(0.12), Inches(1.4), dw - Inches(0.2),
           Inches(1.1), body, sz=8.5, col=LB)

    _r(sl, Inches(0.22), Inches(2.78), Inches(9.56), Inches(0.5), fill=LB, edge=MB)
    _t(sl, Inches(0.35), Inches(2.87), Inches(9.3), Inches(0.38),
       'Pipeline reproductible :   python scripts/run_all_imports.py   →   dbtf build',
       sz=10, bold=True, col=DB, align=PP_ALIGN.CENTER)

    _t(sl, Inches(0.22), Inches(3.4), Inches(9.56), Inches(0.25),
       'QUESTIONS PROBABLES', sz=8, bold=True, col=MB)

    qa = [
        ('Valeurs manquantes ?',
         "Genre → catégorie 'Non renseigné' documentée.  DROM → agrégat DOM."),
        ('Tests dbt ?',
         'Stratégie par couche + test singulier longitudinal + relationships.'),
        ('Reproductibilité ?',
         'Tests bloquants + seeds versionnés + dbtf build suffit.'),
    ]
    for i, (q, a) in enumerate(qa):
        qy = Inches(3.7) + i * Inches(0.6)
        _r(sl, Inches(0.22), qy, Inches(9.56), Inches(0.54), fill=WH, edge=LGR2)
        _t(sl, Inches(0.35), qy + Inches(0.04), Inches(2.4), Inches(0.24),
           f'Q : {q}', sz=8.5, bold=True, col=DB)
        _t(sl, Inches(0.35), qy + Inches(0.28), Inches(9.2), Inches(0.24),
           f'→  {a}', sz=8.5, col=BK)

    _r(sl, 0, H - Inches(0.36), W, Inches(0.36), fill=DB)
    _t(sl, Inches(0.25), H - Inches(0.33), Inches(9.5), Inches(0.3),
       'Yukel Alexandre  ·  Projet 8  ·  Certification Data Analyst OCR  ·  Juin 2026',
       sz=8.5, col=LB, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

BUILDERS = [
    s01_titre, s02_mission, s03_sources, s04_methodo, s05_dag,
    s06_enjeux, s07_tests, s08_regions, s09_genre, s10_insee,
    s11_reco, s12_livrables,
]

def main():
    prs = Presentation(str(V1))
    global W, H
    W, H = prs.slide_width, prs.slide_height
    print(f'Template: {len(prs.slides)} slides  |  {W.inches:.2f}" × {H.inches:.2f}"')

    original = len(prs.slides)
    for fn in BUILDERS:
        fn(prs)
        print(f'  Built: {fn.__name__}')

    # Remove original template slides from manifest
    sldIdLst = prs.slides._sldIdLst
    for _ in range(original):
        sldIdLst.remove(sldIdLst[0])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f'\nSaved → {OUT}')
    print(f'Final: {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
